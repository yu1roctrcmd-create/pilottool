#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
VOL_26 Conflict Zone周辺の飛行 - 改定版PPTXドラフト生成
・Conflict Zone定義・導入、Warship、要撃を削除
・カテゴリー分類・FD Pro表示にフォーカス
"""

from pptx import Presentation
from pptx.util import Emu, Inches
import copy

SRC = "/Users/yuichiromori/Desktop/Conflict Zones周辺の飛行 改定版（2026.04.01）JP.pptx"
DST = "/Users/yuichiromori/Desktop/Conflict Zones周辺の飛行 改定版（2026.04.01）JP_Rev1草案.pptx"

prs = Presentation(SRC)

# =========================================
# 1. 不要スライドを削除（逆順で処理）
#    削除: Slide 3(要撃・Warship), 5(Warship), 6,7,8(要撃シグナル)
#    インデックス: 2, 4, 5, 6, 7
# =========================================
def remove_slide(prs, idx):
    """スライドをインデックス指定で削除"""
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    xml_slides = prs.slides._sldIdLst
    slides = prs.slides
    slide = slides[idx]
    rId = xml_slides[idx].get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
    )
    xml_slides.remove(xml_slides[idx])
    prs.part.drop_rel(rId)

# 逆順で削除（インデックスずれを防ぐ）
for idx in sorted([2, 4, 5, 6, 7], reverse=True):
    remove_slide(prs, idx)

print(f"スライド削除後: {len(prs.slides)}枚")

# =========================================
# 2. Slide 1 の不要シェイプを削除
#    残す: ヘッダー全体・S-3-11紹介balloon・そのパイロット
#    削除: Conflict Zone定義会話4本・OM3-4-2枠・OM3-4-2画像
# =========================================
slide1 = prs.slides[0]

# 削除対象シェイプ名
REMOVE_NAMES_S1 = {
    '円形吹き出し 65',   # "お疲れ！！NCAでは日本-欧州間..."
    'グループ化 69',     # その右パイロット
    '円形吹き出し 42',   # "今、欧州パターンから帰ってきたわ"
    'グループ化 38',     # その左パイロット
    '円形吹き出し 44',   # "このような軍事組織間の紛争地域を..."
    'グループ化 51',     # その右パイロット
    '円形吹き出し 46',   # "そうなんや。我々パイロットの運航にも..."
    'グループ化 56',     # その右パイロット
    '表 2',             # OM3-4-2ヘッダー枠
    '表 9',             # OM3-4-2枠（別）
    'テキスト ボックス 7', # "OM 3-4-2"ラベル
    '6447EBC2-9041-4E5E-BF13-2BEECBC469FC',  # OM3-4-2 内容画像
}

# シェイプを収集してから削除
to_remove = [s for s in slide1.shapes if s.name in REMOVE_NAMES_S1]
for shape in to_remove:
    sp = shape._element
    sp.getparent().remove(sp)
    print(f"  削除: {shape.name}")

# =========================================
# 3. S-3-11紹介balloonを上方向に移動
#    "円形吹き出し 5" と "グループ化 45" を上に移動
# =========================================
MOVE_NAMES = {'円形吹き出し 5', 'グループ化 45'}
SHIFT_UP = Inches(10.0)  # 10インチ上に移動

for shape in slide1.shapes:
    if shape.name in MOVE_NAMES:
        new_top = shape.top - SHIFT_UP
        shape.top = max(new_top, Inches(3.8))  # 最低3.8インチ以上
        print(f"  移動: {shape.name} → top={shape.top/914400:.2f}\"")

# =========================================
# 4. Slide 4（参考文献）のテキストを更新
#    要撃・Warship関連の参考文献を削除
# =========================================
slide4 = prs.slides[2]  # 削除後のインデックス（元Slide4 → now index2）

# 参考文献スライドのテキストボックスを更新
REF_TEXT_NEW = """〇参考文献

OM 3-4 飛行実施計画 3-4-2 ⑦ Conflict Zone
OM S-3-11 Conflict Zoneおよびその周囲での運航について（Eff. 20260401）★NEW
OR OMR　２．Conflict Zone
運航関連業務管理規則　Conflict Zoneに係る飛行空域運用要領
Operation Update No.26003　Explanation of OM S-3-11 'Operations over and near Conflict Zones'　★NEW
Operation Update No.26005　Operations of Flights Considering the Situation in the Middle East　★NEW
IFALPA Flying into and Over Conflict Zones"""

for shape in slide4.shapes:
    if shape.has_text_frame and '参考文献' in shape.text_frame.text:
        tf = shape.text_frame
        # 既存パラグラフをクリアして新テキストを設定
        # まず全パラグラフを取得してフォント情報を保存
        first_para = tf.paragraphs[0]
        # XMLレベルで全パラグラフを削除し新しいテキストを設定
        from pptx.oxml.ns import qn
        from lxml import etree
        txBody = tf._txBody

        # 既存の全 <a:p> を削除
        for p in txBody.findall(qn('a:p')):
            txBody.remove(p)

        # 新しい行を追加
        lines = REF_TEXT_NEW.strip().split('\n')
        for i, line in enumerate(lines):
            p_elem = etree.SubElement(txBody, qn('a:p'))
            pPr = etree.SubElement(p_elem, qn('a:pPr'))

            r_elem = etree.SubElement(p_elem, qn('a:r'))
            rPr = etree.SubElement(r_elem, qn('a:rPr'), attrib={'lang': 'ja-JP', 'altLang': 'en-US'})

            # タイトル行は太字
            if '参考文献' in line or '★NEW' in line:
                rPr.set('b', '1')
            if '★NEW' in line:
                # 赤色
                solidFill = etree.SubElement(rPr, qn('a:solidFill'))
                srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'), attrib={'val': 'CC0000'})

            t_elem = etree.SubElement(r_elem, qn('a:t'))
            t_elem.text = line

        print(f"  参考文献テキスト更新完了")
        break

# =========================================
# 5. 保存
# =========================================
prs.save(DST)
print(f"\n保存完了: {DST}")
print(f"最終スライド数: {len(prs.slides)}枚")
for i, slide in enumerate(prs.slides):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            first = shape.text_frame.paragraphs[0].text[:30].strip()
            if first:
                texts.append(first)
    print(f"  Slide {i+1}: {' | '.join(texts[:3])}")
