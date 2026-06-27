#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Vol.31 Conflict Zones 日本語 → 英語版 変換スクリプト
テキストを英語に置換しつつ、元のデザイン・フォーマットを保持する
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree
import copy

SRC = "/Users/yuichiromori/Desktop/Vol.31 (新)Conflict Zones周辺の飛行_R1.pptx"
DST = "/Users/yuichiromori/Desktop/Vol.31 (New) Flying in Conflict Zones_R1_EN.pptx"

prs = Presentation(SRC)

# =========================================
# 翻訳辞書（形状名 → 英語テキスト）
# キー: (スライドインデックス, 形状名)
# 値: [段落1テキスト, 段落2テキスト, ...]  ※Noneは段落区切り
# =========================================

# 吹き出し内で複数 shape が同名の場合は出現順に対応するため
# 同名の場合は (slide_idx, shape_name, occurrence_index) を使う

TRANSLATIONS = {

    # ── Slide 1 ──────────────────────────────────────────────
    (0, "テキスト ボックス 13"): [
        "NCA Pilot Talk"
    ],
    (0, "テキスト ボックス 18"): [
        "(New)  Flying In and Around Conflict Zones"
    ],

    # Pilot B: お疲れ！！
    (0, "円形吹き出し 65"): [
        "Hey! NCA's Japan-Europe routes have continued to be rerouted through Central Asia.",
        "By the way, the Conflict Zone-related manuals were revised from April 2026, weren't they?"
    ],

    # Pilot C: 欧州帰り（3 paragraphs → merge into logical flow）
    # 同名形状が複数あるので occurrence で分ける → 後述の処理で対応
    # ここでは最初の 円形吹き出し 42 のみ
    (0, "円形吹き出し 44"): [
        "From April 2026, OM S-3-11 was revised, significantly changing how Conflict Zones are handled.",
        "Previously, we basically avoided them at the flight planning stage based on evaluations by public authorities (FAA/EASA, etc.).",
        "Now the company conducts its own risk assessment and classifies Conflict Zones into Categories 1–3 for operational management."
    ],

    (0, "円形吹き出し 26"): [
        "I see. What changes with each category?"
    ],

    # ── Slide 2 ──────────────────────────────────────────────
    (1, "テキスト ボックス 8"): [
        "※ Category determination is based on risk assessment using ICAO Doc.10084.",
        "　(Refer to: Operations Management Regulations / Operation Update No.26005)"
    ],

    # ── Slide 3 ──────────────────────────────────────────────
    (2, "テキスト ボックス 31"): [
        "※ Final determination also considers social/geopolitical context, insurance applicability,",
        "　and other mitigating factors beyond the T/C/V total."
    ],

    # T/C/V タイトルバー
    (2, "Rectangle 21"): [
        "T  Threat — Military/technical risks present in the airspace concerned"
    ],
    (2, "Rectangle 25"): [
        "C  Consequence — Magnitude of impact on the aircraft and aviation environment"
    ],
    (2, "Rectangle 29"): [
        "V  Vulnerability — Status of risk mitigation measures in place"
    ],

    # OM reference labels (Slide 4)
    (3, "テキスト ボックス 19"): [
        "OM S-3-11"
    ],
    (3, "テキスト ボックス 20"): [
        "OM 10-5-2  Distress Communication"
    ],

    # ── Slide 5 ──────────────────────────────────────────────
    # 参考文献はファイル名をそのまま保持しつつ見出しのみ英語化
}

# 吹き出し（同名が複数ある）の翻訳：出現順ごとに定義
# フォーマット: (slide_idx, shape_name) → [出現0の訳, 出現1の訳, ...]
# 訳は「行ごとのリスト」
MULTI_TRANSLATIONS = {

    # Slide 1 ─────────────────────────────────────────────────

    # 円形吹き出し 42  出現0(Pilot C): 欧州帰りの雑談（3パラグラフをまとめる）
    (0, "円形吹き出し 42", 0): [
        "Just got back from the Europe pattern.",
        "Flying along, I thought — the Russia-Ukraine war still hasn't ended,",
        "and conflicts in the Middle East continue. The world is in chaos."
    ],

    # 円形吹き出し 46  出現0(Pilot B): えっそうなんや
    (0, "円形吹き出し 46", 0): [
        "Oh really?! 💦",
        "This really affects our daily operations —",
        "we'd better get our knowledge sorted out!"
    ],

    # TABLE 1x1 (改定概要)  出現1(Slide 1 の2つ目の 1x1 テーブル)
    # ※ テーブルセルで別途処理

    # Slide 2 ─────────────────────────────────────────────────

    # 円形吹き出し 65: カテゴリー分類の評価対象
    (1, "円形吹き出し 65", 0): [
        "The scope of category classification applies to Conflict Zones within 50 NM of the planned flight route,",
        "taking into account the possibility of entry due to deviations.",
        "Crew and dispatchers are required to confirm information on Conflict Zones within at least this range",
        "and provide necessary information sharing and operational support.",
        "In some cases, a change of destination (ATB or DVT) may be required."
    ],

    # 円形吹き出し 5  出現0: 飛行計画段階と飛行中で制限が変わる
    (1, "円形吹き出し 5", 0): [
        "The restrictions differ between the flight planning stage and during flight. Here's a summary:"
    ],

    # 円形吹き出し 26  出現0: 飛行中も慎重なモニター
    (1, "円形吹き出し 26", 0): [
        "Careful monitoring during flight is essential too.",
        "By the way, is it true that Conflict Zones can now be seen on FD Pro?"
    ],

    # 円形吹き出し 5  出現1: FD Pro の説明
    (1, "円形吹き出し 5", 1): [
        "Yes! They're displayed color-coded by category on the Enroute Chart.",
        "Tap the Boundary Line or Ball Note to view detailed information."
    ],

    # 円形吹き出し 26  出現1: Data is current
    (1, "円形吹き出し 26", 1): [
        "It's great to be able to visually check during pre-flight briefing and during flight.",
        "Since the display needs to be updated whenever new information is available,",
        "make sure to confirm \"Data is current\" before departure."
    ],

    # Slide 3 ─────────────────────────────────────────────────

    # 円形吹き出し 65: 茶本の説明 + リスク評価
    (2, "円形吹き出し 65", 0): [
        "The Operations Management Regulations (the 'brown book') covers not only Conflict Zones",
        "but also various flight operations and safety-related rules — well worth reading.",
        "Now, regarding the risk assessment:",
        "It analyzes three factors: Threat, Consequence, and Vulnerability."
    ],

    # 円形吹き出し 46: カテゴリーはどうやって決まるの？
    (2, "円形吹き出し 46", 0): [
        "By the way, how is the category determined?",
        "There's a note saying it's 'regulated in the Operations Management Regulations,'",
        "but honestly, I've never actually looked at it..."
    ],

    # 円形吹き出し 5: 詳しい説明（Rate 1-3）
    (2, "円形吹き出し 5", 0): [
        "Let me explain in a bit more detail.",
        "Each of the 3 factors is rated from Rate 1 to 3,",
        "and the total score determines the risk level:",
        "  High → Category 1",
        "  Medium → Category 2",
        "  Low → Category 3"
    ],

    # Rectangle 23: T の Rate 説明
    (2, "Rectangle 23", 0): [
        "Rate 3 (High)",
        "Actual military action/attack has occurred in recent years, OR clear evidence",
        "of intent, capability, or plan to attack exists.",
        "Rate 2 (Medium)",
        "Relatively recent evidence or reports indicate a short/medium-term attack plan or intent.",
        "Rate 1 (Low)",
        "No recent incidents; no signs of attacks or attack planning."
    ],

    # Rectangle 27: C の Rate 説明
    (2, "Rectangle 27", 0): [
        "Rate 3 (Severe)",
        "Loss of aircraft, OR serious disruption to ATC/traffic flow (e.g., paralysis).",
        "Rate 2 (Moderate)",
        "Serious aircraft damage, OR significant impact on ATC/traffic flow",
        "(e.g., major delays to traffic passage).",
        "Rate 1 (Minor)",
        "Partial aircraft damage, OR limited impact on ATC/traffic flow."
    ],

    # Rectangle 31: V の Rate 説明
    (2, "Rectangle 31", 0): [
        "Rate 3 (High)",
        "No risk mitigation measures (e.g., flight restrictions to avoid risk) have been implemented.",
        "Rate 2 (Medium)",
        "Mitigation measures exist but are currently insufficient or difficult to apply effectively.",
        "Rate 1 (Low)",
        "Effective and recognized risk mitigation measures are in place."
    ],

    # Slide 4 ─────────────────────────────────────────────────

    # 円形吹き出し 65: Cさんありがとう
    (3, "円形吹き出し 65", 0): [
        "Thanks, C! We should also review the procedures for unexpected situations",
        "during flight — specifically OM S-3-11 and 10-5-2."
    ],

    # 円形吹き出し 46  出現0: みんな
    (3, "円形吹き出し 46", 0): [
        "Everyone's really studying hard!",
        "Have there been any changes to the OM reporting requirements?"
    ],

    # 円形吹き出し 46  出現1: 家に帰って + ジョーク
    (3, "円形吹き出し 46", 1): [
        "When I get home today I'll review all the changes after some rest.",
        "Oh wait — I had an argument with my wife before leaving,",
        "and the home situation is already 'Category 1'...",
        "If intercepted, I'll squawk 7700 and comply with her instructions — yes ma'am!! 😅"
    ],

    # 円形吹き出し 5: ASR の説明
    (3, "円形吹き出し 5", 0): [
        "Great question! Conflict Zone-related events have been added to the ASR submission requirements.",
        "Reports from flight crew are one important source of information for category determination,",
        "so please don't forget to report!"
    ],

    # Rectangle 23: ASR 要件ボックス
    (3, "Rectangle 23", 0): [
        "OM S-3-11  Section 3.4  Post-Flight",
        "The Captain shall submit an ASR in the following cases:",
        "・ Entered, or considered entering, a Conflict Zone to avoid adverse weather or for other reasons",
        "・ GPS Jamming / Spoofing occurred or was suspected",
        "・ Any other case the Captain deems necessary regarding Conflict Zone operations"
    ],
}

# テーブルセルの翻訳
# (slide_idx, shape_name, row, col) → 英語テキスト
TABLE_TRANSLATIONS = {

    # Slide 1 ─ 改定概要テーブル (TABLE 1x1, 出現1)
    # 出現0は空白テーブルなので出現1が改定概要
    (0, "TABLE_OCC1", 0, 0): (
        "【Revision Summary】\n"
        "OM S-3-11 effective 2026/04/01 established a new classification system for Conflict Zones.\n"
        "This Vol.31 is an update reflecting those changes.\n"
        "For details on Interception procedures, please also refer to the previous article Vol.26 "
        "\"Flying In and Around Conflict Zones.\""
    ),

    # Slide 2 ─ カテゴリーテーブル (TABLE 4x4)
    (1, "TABLE 4x4", 0, 0): "Category",
    (1, "TABLE 4x4", 0, 1): "Classification",
    (1, "TABLE 4x4", 0, 2): "Flight Planning Stage",
    (1, "TABLE 4x4", 0, 3): "In Flight",

    (1, "TABLE 4x4", 1, 0): "Category 1",
    (1, "TABLE 4x4", 1, 1): "Prohibited",
    (1, "TABLE 4x4", 1, 2): "Avoid the area",
    (1, "TABLE 4x4", 1, 3): "Entry into the area is prohibited",

    (1, "TABLE 4x4", 2, 0): "Category 2",
    (1, "TABLE 4x4", 2, 1): "Restricted",
    (1, "TABLE 4x4", 2, 2): (
        "Avoid the specified altitude(s) within the area,\n"
        "unless flight through the area at those altitudes is operationally unavoidable\n"
        "and separately specified requirements are satisfied."
    ),
    (1, "TABLE 4x4", 2, 3): (
        "Avoid the specified altitude(s) within the area,\n"
        "unless the requirements confirmed during the flight planning stage are satisfied."
    ),

    (1, "TABLE 4x4", 3, 0): "Category 3",
    (1, "TABLE 4x4", 3, 1): "Information",
    (1, "TABLE 4x4", 3, 2): "Normal operations are permitted, maintaining awareness of the latest available information.",
    (1, "TABLE 4x4", 3, 3): "",

    # Slide 2 ─ FD Pro テーブル (TABLE 4x2)
    (1, "TABLE 4x2", 0, 0): "Category",
    (1, "TABLE 4x2", 0, 1): "FD Pro Display",
    (1, "TABLE 4x2", 1, 0): "Category 1",
    (1, "TABLE 4x2", 1, 1): "Brown dotted line + fill",
    (1, "TABLE 4x2", 2, 0): "Category 2",
    (1, "TABLE 4x2", 2, 1): "Red dotted line + fill",
    (1, "TABLE 4x2", 3, 0): "Category 3",
    (1, "TABLE 4x2", 3, 1): "Brown dotted line only",

    # Slide 3 ─ リスクレベル → カテゴリー表 (TABLE 4x5)
    (2, "TABLE 4x5", 0, 0): "Risk Level",
    (2, "TABLE 4x5", 0, 1): "T+C+V",
    (2, "TABLE 4x5", 0, 2): "Category",
    (2, "TABLE 4x5", 0, 3): "Restriction",
    (2, "TABLE 4x5", 0, 4): "Operational Condition",

    (2, "TABLE 4x5", 1, 0): "High",
    (2, "TABLE 4x5", 1, 1): "7–9",
    (2, "TABLE 4x5", 1, 2): "Category 1\n(Prohibited)",
    (2, "TABLE 4x5", 1, 3): "Flight prohibited",
    (2, "TABLE 4x5", 1, 4): "Risk mitigation measures\nmust be implemented.",

    (2, "TABLE 4x5", 2, 0): "Medium",
    (2, "TABLE 4x5", 2, 1): "4–6",
    (2, "TABLE 4x5", 2, 2): "Category 2\n(Restricted)",
    (2, "TABLE 4x5", 2, 3): "Flight restricted",
    (2, "TABLE 4x5", 2, 4): "Confirm that existing measures\nare in place (considering\nfactors beyond the record).",

    (2, "TABLE 4x5", 3, 0): "Low",
    (2, "TABLE 4x5", 3, 1): "1–3",
    (2, "TABLE 4x5", 3, 2): "Category 3\n(Information)",
    (2, "TABLE 4x5", 3, 3): "Normal ops",
    (2, "TABLE 4x5", 3, 4): "No specific measures required.",
}

# =========================================
# テキスト置換ユーティリティ
# =========================================

def set_shape_text(shape, lines):
    """
    shape のテキストフレームを lines（リスト）で置き換える。
    元のフォーマット（フォントサイズ・色など）を最初のランから継承する。
    """
    if not shape.has_text_frame:
        return
    tf = shape.text_frame

    # 元の書式を取得（最初のパラグラフ・最初のラン）
    orig_para = tf.paragraphs[0] if tf.paragraphs else None
    orig_run  = orig_para.runs[0] if (orig_para and orig_para.runs) else None

    def get_font_props(run):
        props = {}
        if run:
            try: props['size']   = run.font.size
            except: pass
            try: props['bold']   = run.font.bold
            except: pass
            try: props['italic'] = run.font.italic
            except: pass
            try: props['color']  = run.font.color.rgb
            except: pass
        return props

    orig_props = get_font_props(orig_run)

    # 既存コンテンツをクリア（XML レベルで全 <a:p> 削除後再構築）
    txBody = tf._txBody
    for p in txBody.findall(qn('a:p')):
        txBody.remove(p)

    for line in lines:
        p_elem = etree.SubElement(txBody, qn('a:p'))
        r_elem = etree.SubElement(p_elem, qn('a:r'))
        rPr    = etree.SubElement(r_elem, qn('a:rPr'),
                                  attrib={'lang': 'en-US', 'altLang': 'en-US'})
        if orig_props.get('size'):
            rPr.set('sz', str(int(orig_props['size'] / 127)))  # EMU→hundredths of pt
        if orig_props.get('bold'):
            rPr.set('b', '1')
        t_elem = etree.SubElement(r_elem, qn('a:t'))
        t_elem.text = line


def set_table_cell_text(cell, text):
    """テーブルセルのテキストを置き換える（改行対応）"""
    tf = cell.text_frame

    # 元の書式を取得
    orig_run = None
    if tf.paragraphs and tf.paragraphs[0].runs:
        orig_run = tf.paragraphs[0].runs[0]

    orig_size  = None
    orig_bold  = None
    orig_color = None
    if orig_run:
        try: orig_size  = orig_run.font.size
        except: pass
        try: orig_bold  = orig_run.font.bold
        except: pass
        try: orig_color = orig_run.font.color.rgb
        except: pass

    txBody = tf._txBody
    for p in txBody.findall(qn('a:p')):
        txBody.remove(p)

    for line in text.split('\n'):
        p_elem = etree.SubElement(txBody, qn('a:p'))
        r_elem = etree.SubElement(p_elem, qn('a:r'))
        rPr    = etree.SubElement(r_elem, qn('a:rPr'),
                                  attrib={'lang': 'en-US', 'altLang': 'en-US'})
        if orig_size:
            rPr.set('sz', str(int(orig_size / 127)))
        if orig_bold:
            rPr.set('b', '1')
        if orig_color:
            sf = etree.SubElement(rPr, qn('a:solidFill'))
            cl = etree.SubElement(sf,  qn('a:srgbClr'))
            cl.set('val', str(orig_color))
        t_elem = etree.SubElement(r_elem, qn('a:t'))
        t_elem.text = line


# =========================================
# メイン処理
# =========================================

for slide_idx, slide in enumerate(prs.slides):

    # ── 形状ごとの出現カウンター（同名形状の区別に使用）
    shape_occurrence = {}

    # ── テーブル識別用カウンター（サイズで区別）
    table_size_count = {}

    for shape in slide.shapes:
        name = shape.name

        # ─ 出現インデックスを追跡
        occ = shape_occurrence.get(name, 0)
        shape_occurrence[name] = occ + 1

        # ─ SINGLE_KEY で翻訳があれば置換（最初の出現のみ）
        single_key = (slide_idx, name)
        if single_key in TRANSLATIONS and occ == 0:
            set_shape_text(shape, TRANSLATIONS[single_key])
            continue

        # ─ MULTI_KEY で翻訳があれば置換
        multi_key = (slide_idx, name, occ)
        if multi_key in MULTI_TRANSLATIONS:
            set_shape_text(shape, MULTI_TRANSLATIONS[multi_key])
            continue

        # ─ テーブル処理
        if shape.shape_type == 19:  # TABLE
            tbl = shape.table
            nr  = len(tbl.rows)
            nc  = len(tbl.columns)
            size_key = f"TABLE {nr}x{nc}"
            size_occ = table_size_count.get(size_key, 0)
            table_size_count[size_key] = size_occ + 1

            # 改定概要テーブル (1x1, 出現1)
            if size_key == "TABLE 1x1" and size_occ == 1:
                txt = TABLE_TRANSLATIONS.get((slide_idx, "TABLE_OCC1", 0, 0), "")
                if txt:
                    set_table_cell_text(tbl.cell(0, 0), txt)
                continue

            # 通常テーブル
            for ri in range(nr):
                for ci in range(nc):
                    tbl_key = (slide_idx, size_key, ri, ci)
                    if tbl_key in TABLE_TRANSLATIONS:
                        set_table_cell_text(tbl.cell(ri, ci),
                                            TABLE_TRANSLATIONS[tbl_key])

# ── Slide 5: 参考文献の見出しを英語化
slide5 = prs.slides[4]
for shape in slide5.shapes:
    if shape.has_text_frame and '参考文献' in shape.text_frame.text:
        tf = shape.text_frame
        txBody = tf._txBody
        for p in txBody.findall(qn('a:p')):
            txBody.remove(p)
        refs = [
            ("○ References", True),
            ("", False),
            ("OM 3-4 Flight Planning  3-4-2 ⑦ Conflict Zone", False),
            ("OM S-3-11  Operations over and near Conflict Zones", False),
            ("OR OMR  2. Conflict Zone", False),
            ("Operations Management Regulations  —  Airspace Operations Procedures for Conflict Zones", False),
            ("Operation Update No.26003  Explanation of OM S-3-11 'Operations over and near Conflict Zones'", False),
            ("Operation Update No.26005  Operations of Flights Considering the Situation in the Middle East", False),
        ]
        for text, bold in refs:
            p_elem = etree.SubElement(txBody, qn('a:p'))
            r_elem = etree.SubElement(p_elem, qn('a:r'))
            rPr    = etree.SubElement(r_elem, qn('a:rPr'),
                                      attrib={'lang': 'en-US', 'altLang': 'en-US'})
            if bold:
                rPr.set('b', '1')
            t_elem = etree.SubElement(r_elem, qn('a:t'))
            t_elem.text = text
        break

# =========================================
# 保存
# =========================================
prs.save(DST)
print(f"保存完了: {DST}")
print("\n各スライドのテキスト確認:")
prs2 = Presentation(DST)
for i, slide in enumerate(prs2.slides):
    print(f"\n=== Slide {i+1} ===")
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    print(f"  {t[:80]}")
        if shape.shape_type == 19:
            for r in shape.table.rows:
                row_t = [c.text.strip()[:25] for c in r.cells]
                if any(row_t):
                    print(f"  TABLE ROW: {row_t}")
