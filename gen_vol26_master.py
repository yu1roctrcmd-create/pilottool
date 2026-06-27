#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
VOL_26 Conflict Zone 改定版 ─ マスタースクリプト（最終版）

デザイン基準（Slide 1・2 に合わせる）
  吹き出し色: 緑=#33CC33 / 黄=#FFC000 / 青=#00B0F0
  本文フォント: 14pt（Slide 2 と統一）
  ヘッダー: Slide 1 から構造ごとコピー

会話の流れ:
  Slide 1: 欧州帰りの雑談 → OM S-3-11 改訂・カテゴリー分類の紹介
  Slide 2: カテゴリーの詳細・50NM基準・FD Pro表示 → 評価方法への疑問
  Slide 3: T×C×V リスク評価の仕組み → ASRへの話題転換
  Slide 4: ASR 報告義務・不測の事態の対処 → 笑いで締め
  Slide 5: 参考文献
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.oxml.ns import qn
from lxml import etree
import copy

SRC = "/Users/yuichiromori/Desktop/Conflict Zones周辺の飛行 改定版（2026.04.01）JP.pptx"
DST = "/Users/yuichiromori/Desktop/Conflict Zones周辺の飛行 改定版（2026.04.01）JP_Rev1草案.pptx"

prs = Presentation(SRC)

# ─────────────── カラー定数（Slide 1・2 から抽出）───────────────
C_GREEN  = (0x33, 0xCC, 0x33)   # Pilot C の吹き出し
C_YELLOW = (0xFF, 0xC0, 0x00)   # Pilot B の吹き出し
C_BLUE   = (0x00, 0xB0, 0xF0)   # Pilot A の吹き出し（説明担当）
C_WHITE  = (255, 255, 255)
C_DARK   = (30, 30, 30)
C_GRAY   = (100, 100, 100)
C_NAVY   = (26, 54, 93)

# ─────────────── ユーティリティ ───────────────
def table_cell_fill(cell, hex_color):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for ex in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(ex)
    sf  = etree.SubElement(tcPr, qn('a:solidFill'))
    cl  = etree.SubElement(sf,   qn('a:srgbClr'))
    cl.set('val', hex_color)

def set_tbl_cell(tbl, row, col, text, fill=None, bold=False, size=9,
                 fc=(30,30,30), align=PP_ALIGN.LEFT):
    cell = tbl.cell(row, col)
    cell.text = text
    if fill:
        table_cell_fill(cell, fill)
    for para in cell.text_frame.paragraphs:
        para.alignment = align
        for run in para.runs:
            run.font.size  = Pt(size)
            run.font.bold  = bold
            run.font.color.rgb = RGBColor(*fc)

def add_tb(slide, text, l, t, w, h, size=12, bold=False,
           color=C_DARK, align=PP_ALIGN.LEFT,
           bg=None, border=None):
    """シンプルなテキストボックス"""
    tb = slide.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left   = Inches(0.08)
    tf.margin_top    = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.margin_right  = Inches(0.08)
    if bg:
        tb.fill.solid()
        tb.fill.fore_color.rgb = RGBColor(*bg)
    if border:
        tb.line.color.rgb = RGBColor(*border)
        tb.line.width = Pt(0.75)
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size  = Pt(size)
        r.font.bold  = bold
        r.font.color.rgb = RGBColor(*color)
    return tb

def balloon(slide, speaker, text, l, t, w, h,
            fill=C_YELLOW, size=14):
    """
    吹き出し風テキストボックス（Slide 1・2 のスタイルに合わせる）
    - fill: 吹き出し背景色 (RGB tuple)
    - size: 本文フォントサイズ（デフォルト14pt）
    - スピーカー名は小さいイタリックで先頭に付与
    """
    tb = slide.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left   = Inches(0.15)
    tf.margin_top    = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    tf.margin_right  = Inches(0.15)
    tb.fill.solid()
    tb.fill.fore_color.rgb = RGBColor(*fill)
    # 枠線は同色より少し暗く
    dark = tuple(max(0,c-60) for c in fill)
    tb.line.color.rgb = RGBColor(*dark)
    tb.line.width = Pt(0.75)

    # ─ スピーカー名（小さめ・イタリック）
    p0 = tf.paragraphs[0]
    rs = p0.add_run()
    rs.text     = f"（{speaker}）"
    rs.font.size   = Pt(size - 3)
    rs.font.italic = True
    rs.font.bold   = False
    # 吹き出し色が明るければ暗い文字、暗ければ白文字
    brightness = sum(fill) / 3
    rs.font.color.rgb = RGBColor(50,50,50) if brightness > 150 else RGBColor(220,220,220)

    # ─ 本文
    for line in text.split('\n'):
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.size  = Pt(size)
        r.font.color.rgb = RGBColor(30,30,30) if brightness > 150 else RGBColor(255,255,255)
    return tb

def copy_header_from_slide1(dst_slide, src_slide, new_title):
    """
    Slide 1 のヘッダー要素（テキストのみ。図はスキップ）を
    dst_slide にコピーし、タイトルテキストを new_title に差し替える
    """
    HEADER_SHAPE_NAMES = {
        '正方形/長方形 8',      # 濃紺背景
        'テキスト ボックス 10',  # SHIBAYAMA SMALL TOWN TALK
        'テキスト ボックス 13',  # NCA 井戸端会議
        'テキスト ボックス 15',  # Vol. 26 ...
        '直線コネクタ 16',       # 区切り線
    }
    spTree = dst_slide.shapes._spTree
    for shape in src_slide.shapes:
        if shape.name in HEADER_SHAPE_NAMES:
            spTree.append(copy.deepcopy(shape._element))

    # タイトル帯（白背景）と新しいタイトルテキスト
    tbg = dst_slide.shapes.add_shape(
        1, Inches(0.39), Inches(2.62), Inches(12.26), Inches(0.72))
    tbg.fill.solid()
    tbg.fill.fore_color.rgb = RGBColor(255,255,255)
    tbg.line.fill.background()

    tb_t = dst_slide.shapes.add_textbox(
        Inches(0.81), Inches(2.65), Inches(11.5), Inches(0.6))
    p_t = tb_t.text_frame.paragraphs[0]
    p_t.alignment = PP_ALIGN.CENTER
    r_t = p_t.add_run()
    r_t.text = new_title
    r_t.font.size  = Pt(16)
    r_t.font.bold  = True
    r_t.font.color.rgb = RGBColor(*C_NAVY)

# =========================================
# PHASE 1: 不要スライドを削除
#   削除: Slide 3(要撃・Warship), 5(Warship詳細), 6,7,8(要撃シグナル)
# =========================================
def remove_slide(prs, idx):
    sldIdLst = prs.slides._sldIdLst
    rId = sldIdLst[idx].get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    sldIdLst.remove(sldIdLst[idx])
    prs.part.drop_rel(rId)

for idx in sorted([2,4,5,6,7], reverse=True):
    remove_slide(prs, idx)
print(f"スライド削除後: {len(prs.slides)}枚")

# =========================================
# PHASE 2: Slide 1 の修正
#   - Conflict Zone 定義の吹き出し・OM3-4-2 を削除
#   - S-3-11 吹き出しを上に移動
#   - 導入会話（Pilot C→B→A）を追加
# =========================================
slide1 = prs.slides[0]

REMOVE_S1 = {
    '円形吹き出し 65','グループ化 69',
    '円形吹き出し 42','グループ化 38',
    '円形吹き出し 44','グループ化 51',
    '円形吹き出し 46','グループ化 56',
    '表 2','表 9','テキスト ボックス 7',
    '6447EBC2-9041-4E5E-BF13-2BEECBC469FC',
}
for s in [s for s in slide1.shapes if s.name in REMOVE_S1]:
    s._element.getparent().remove(s._element)

# S-3-11 balloon と隣のパイロットを上に移動
SHIFT = Inches(10.0)
for s in slide1.shapes:
    if s.name in {'円形吹き出し 5', 'グループ化 45'}:
        s.top = max(s.top - SHIFT, Inches(7.2))

# ─ 導入会話（S-3-11 balloon の前）
balloon(slide1, "Pilot C",
    "欧州パターンから帰ってきたわ。\n"
    "ロシア・ウクライナ戦争もまだ終わらへんし、中東の紛争もある。\n"
    "世界は混沌としてるな〜。",
    2.4, 3.65, 8.8, 1.2, fill=C_GREEN)

balloon(slide1, "Pilot B",
    "NCAでも迂回ルートが続いてるよね。\n"
    "そういえば今年の4月から Conflict Zone まわりの規程が変わったって聞いたよ？",
    1.2, 5.0, 9.0, 1.05, fill=C_YELLOW)

# ─ 既存 S-3-11 balloon の上にスピーカーラベルを追加
add_tb(slide1, "（Pilot A）", 1.2, 7.1, 2.0, 0.35,
       size=11, bold=False, color=C_GRAY)

# =========================================
# PHASE 3: Slide 2 末尾に橋渡し会話を追加
#   （カテゴリー表・FD Pro → リスク評価の疑問）
# =========================================
slide2 = prs.slides[1]

balloon(slide2, "Pilot B",
    "ところで、カテゴリーはどうやって決まるの？\n"
    "※に「ICAO Doc.10084 をベースにしたリスク評価」って書いてあるけど、\n"
    "具体的にはどんな評価をするの？",
    2.5, 16.7, 8.8, 1.3, fill=C_YELLOW)

balloon(slide2, "Pilot A",
    "T（脅威）、C（結果の重大性）、V（脆弱性）の\n"
    "3要素をそれぞれ評価するんだよ。次のスライドで詳しく説明するね！",
    1.2, 18.15, 9.0, 1.0, fill=C_BLUE)

# =========================================
# PHASE 4: Slide 3（新規）─ リスク評価方法
# =========================================
blank = prs.slide_layouts[6]
slide3 = prs.slides.add_slide(blank)
copy_header_from_slide1(slide3, slide1,
    "カテゴリー評価方法 ─ T × C × V リスクアセスメント")

# ─ 導入会話
balloon(slide3, "Pilot B",
    "T・C・V ってそれぞれ何を評価するの？",
    2.5, 3.55, 8.8, 0.75, fill=C_YELLOW)

balloon(slide3, "Pilot A",
    "それぞれを Rate 1〜3 で評価して、合計値（TTL Rate）で\n"
    "リスクレベルを決めるんだ。\n"
    "High→Category 1、Medium→Category 2、Low→Category 3 に対応するよ。",
    1.2, 4.45, 10.5, 1.2, fill=C_BLUE)

# ─ T/C/V 3要素ブロック
factors = [
    ("T  Threat（脅威）─ 当該空域に存在する軍事的・技術的な危険性",
     ("Rate 3（高い）",
      "　過去数年以内に実際の軍事行動・攻撃が発生、またはその意図・能力の明確な証拠がある。\n"
      "Rate 2（中程度）\n"
      "　比較的最近に、短期・中期の攻撃計画や意図を示す事例・証拠がある。\n"
      "Rate 1（低い）\n"
      "　最近の事例がなく、攻撃や攻撃計画等の兆候もない。"),
     "CC3333","FFD7D7"),
    ("C  Consequence（結果の重大性）─ 機体および航空環境への影響の大きさ",
     ("Rate 3（重大）",
      "　機体の損失、または管制・交通流への深刻な影響（交通流の麻痺等）。\n"
      "Rate 2（中程度）\n"
      "　機体の重篤な損傷、または管制・交通流への大きな影響（通過の大幅な遅延等）。\n"
      "Rate 1（軽微）\n"
      "　機体の部分的な損傷、または管制・交通流への部分的な影響。"),
     "CC9900","FFF2CC"),
    ("V  Vulnerability（脆弱性）─ リスク緩和措置の実施状況",
     ("Rate 3（高い）",
      "　リスクの緩和措置（リスク回避のための制限等）が実施されていない。\n"
      "Rate 2（中程度）\n"
      "　緩和措置は存在するが、現時点で効果的な対策の設定が困難。\n"
      "Rate 1（低い）\n"
      "　効果的と認められるリスクの緩和措置が実施されている。"),
     "336633","E2EFDA"),
]

y = 5.8
for title, (r3_label, body_text), title_hex, body_hex in factors:
    # タイトル帯
    r_,g_,b_ = int(title_hex[0:2],16), int(title_hex[2:4],16), int(title_hex[4:6],16)
    tbg = slide3.shapes.add_shape(1, Inches(0.8), Inches(y), Inches(11.5), Inches(0.35))
    tbg.fill.solid(); tbg.fill.fore_color.rgb = RGBColor(r_,g_,b_)
    tbg.line.fill.background()
    add_tb(slide3, title, 0.88, y+0.02, 11.3, 0.31,
           size=10, bold=True, color=C_WHITE)
    y += 0.35
    # 本文（Rate 3 ラベルを太字で先頭に）
    br_,bg_,bb_ = int(body_hex[0:2],16), int(body_hex[2:4],16), int(body_hex[4:6],16)
    bbg = slide3.shapes.add_shape(1, Inches(0.8), Inches(y), Inches(11.5), Inches(1.0))
    bbg.fill.solid(); bbg.fill.fore_color.rgb = RGBColor(br_,bg_,bb_)
    bbg.line.color.rgb = RGBColor(r_,g_,b_); bbg.line.width = Pt(0.5)

    tb_body = slide3.shapes.add_textbox(Inches(0.9),Inches(y+0.04),Inches(11.2),Inches(0.92))
    tf = tb_body.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_top = Inches(0.03)
    # Rate 3 ラベル（太字）
    p0 = tf.paragraphs[0]
    r0 = p0.add_run(); r0.text = r3_label
    r0.font.size = Pt(10); r0.font.bold = True
    r0.font.color.rgb = RGBColor(r_,g_,b_)
    # 残りテキスト
    for i, line in enumerate(body_text.split('\n')):
        p = tf.add_paragraph()
        r = p.add_run(); r.text = line
        r.font.size = Pt(9.5)
        # "Rate X" の行は太字
        r.font.bold = line.strip().startswith('Rate')
        r.font.color.rgb = RGBColor(r_,g_,b_) if line.strip().startswith('Rate') else RGBColor(*C_DARK)
    y += 1.05

# ─ リスクレベル → カテゴリー 決定表
add_tb(slide3, "▼ T + C + V 合計値によるカテゴリー決定",
    0.8, y+0.1, 11.5, 0.38, size=11, bold=True, color=C_NAVY)
y += 0.53

tbl = slide3.shapes.add_table(
    4, 5, Inches(0.8), Inches(y), Inches(11.5), Inches(1.6)).table
for ci, w in enumerate([1.2,1.3,2.1,3.0,3.9]):
    tbl.columns[ci].width = Inches(w)
for ci, h in enumerate(["Risk Level","T+C+V合計","Category","制限内容","運航上の条件"]):
    set_tbl_cell(tbl,0,ci,h,"2E5FA3",bold=True,size=10,fc=(255,255,255),align=PP_ALIGN.CENTER)
rows = [
    ("High",  "7〜9","Category 1\n(Prohibited)","飛行禁止",
     "リスク抑制のための対策を\n講じる必要がある","FFD7D7"),
    ("Medium","4〜6","Category 2\n(Restricted)","飛行制限",
     "既にリスク抑制措置が含まれていること\nを確認（記録外要素も考慮）","FFF2CC"),
    ("Low",   "1〜3","Category 3\n(Information)","通常運航可",
     "具体的な措置等は不要","E2EFDA"),
]
for ri,(lvl,ttl,cat,lim,cond,col) in enumerate(rows):
    for ci,txt in enumerate([lvl,ttl,cat,lim,cond]):
        set_tbl_cell(tbl,ri+1,ci,txt,col,size=10,
                     align=PP_ALIGN.CENTER if ci<3 else PP_ALIGN.LEFT)

add_tb(slide3,
    "※ 最終決定では T/C/V 合計に加え、社的状況・保険の適用・その他の緩和要素も勘案。\n"
    "   評価結果は Risk Assessment Record に記録し、状況変化時は再評価を実施。（OR 運航関連業務管理規則 Appendix参照）",
    0.8, y+1.65, 11.5, 0.5, size=9, color=C_GRAY)

# ─ Slide 3 → Slide 4 への橋渡し
balloon(slide3, "Pilot C",
    "なるほど、しっかりした根拠でカテゴリーが決まるんだね。\n"
    "ところでフライト後の報告義務って何か変わったの？",
    2.5, y+2.22, 8.8, 1.0, fill=C_GREEN)

balloon(slide3, "Pilot A",
    "いい質問！OM S-3-11 の新設で\n"
    "ASR の報告対象に Conflict Zone 関連事象が新たに追加されたよ。次で詳しく！",
    1.2, y+3.36, 10.5, 1.0, fill=C_BLUE)

# =========================================
# PHASE 5: Slide 4（新規）─ ASR 報告義務
# =========================================
slide4 = prs.slides.add_slide(blank)
copy_header_from_slide1(slide4, slide1,
    "飛行後の報告 ── ASR 報告対象の新規追加")

balloon(slide4, "Pilot A",
    "OM S-3-11 の新設に伴い、Air Safety Report（ASR）の報告対象に\n"
    "Conflict Zone 関連の事象が新たに追加されたよ。",
    1.2, 3.55, 10.5, 1.0, fill=C_BLUE)

# ─ ASR 対象表
tbl2 = slide4.shapes.add_table(
    4, 3, Inches(0.8), Inches(4.65), Inches(11.5), Inches(1.8)).table
tbl2.columns[0].width = Inches(0.5)
tbl2.columns[1].width = Inches(8.1)
tbl2.columns[2].width = Inches(2.9)
for ci, h in enumerate(["#","ASR 報告の対象となる事象","関連規程"]):
    set_tbl_cell(tbl2,0,ci,h,"2E5FA3",bold=True,size=10,fc=(255,255,255),align=PP_ALIGN.CENTER)
asrs = [
    ("①","悪天回避等により Conflict Zone（計画経路から 50NM 以内）へ\n入域した、または入域を検討した場合","OM S-3-11\nOR 22章"),
    ("②","GPS Jamming / Spoofing が発生した、または疑われた場合","OM S-3-11\nOR 22章"),
    ("③","Conflict Zone 飛行に係るその他の特異事象","OM S-3-11"),
]
fills4 = ["F5F5F5","FFFFFF","F5F5F5"]
for ri,((no,ev,ref),fl) in enumerate(zip(asrs,fills4)):
    set_tbl_cell(tbl2,ri+1,0,no,  fl,size=10,align=PP_ALIGN.CENTER)
    set_tbl_cell(tbl2,ri+1,1,ev,  fl,size=10)
    set_tbl_cell(tbl2,ri+1,2,ref, fl,size=10)

balloon(slide4, "Pilot C",
    "悪天等での Deviation 中に Conflict Zone に入域したり、\n"
    "GPS Jamming / Spoofing が疑われたら ASR を出すんだね。\n"
    "万が一、入域が避けられない場合の対処措置も確認しておかないとね！",
    2.5, 6.6, 8.8, 1.25, fill=C_GREEN)

# ─ 不測の事態の対処措置
add_tb(slide4, "▼ Conflict Zone への入域が避けられない場合の対処措置",
    0.8, 8.0, 11.5, 0.4, size=11, bold=True, color=C_NAVY)

tb_m = slide4.shapes.add_textbox(Inches(0.8),Inches(8.46),Inches(11.5),Inches(1.55))
tf_m = tb_m.text_frame; tf_m.word_wrap = True
tf_m.margin_left = Inches(0.12); tf_m.margin_top = Inches(0.08)
tb_m.fill.solid(); tb_m.fill.fore_color.rgb = RGBColor(255,248,220)
tb_m.line.color.rgb = RGBColor(200,160,60); tb_m.line.width = Pt(0.75)
measures = [
    ("・ トランスポンダ 7700 のセット", False),
    ("・ 遭難通信の発信", False),
    ("・ 121.5 MHz での一方送信", False),
    ("・ 夜間のエクステリアライト点灯（通常 MEL 範囲内では機体識別への影響なし）", False),
    ("※ 緊急事態または安全上やむを得ない入域は直ちに危険をもたらすものではない。安全最優先で対処すること。", True),
]
for i,(line,is_note) in enumerate(measures):
    p = tf_m.paragraphs[0] if i==0 else tf_m.add_paragraph()
    r = p.add_run(); r.text = line
    r.font.size = Pt(11)
    r.font.bold = is_note
    r.font.color.rgb = RGBColor(160,0,0) if is_note else RGBColor(*C_DARK)

add_tb(slide4, "（参照：OM S-3-11 3.5節 / OR 運航関連業務管理規則 22章）",
    0.8, 10.06, 11.5, 0.3, size=9, color=C_GRAY)

# ─ 締めのジョーク（Pilot A）
balloon(slide4, "Pilot A",
    "今日は帰ったら変更点のレビューせんとなぁ。\n"
    "そういや妻とのケンカは Category 1 のままで、\n"
    "Vulnerability は Rate 3（緩和措置なし）… ASR の提出が必要かも！",
    1.2, 10.5, 10.5, 1.2, fill=C_BLUE)

balloon(slide4, "Pilot C",
    "それは Category 1 なら要撃は必至やで（笑）\n"
    "今回の改定をしっかり把握して、安全な飛行を！",
    2.5, 11.85, 8.8, 1.0, fill=C_GREEN)

# =========================================
# PHASE 6: Slide 5（参考文献）更新
# =========================================
slide_ref = prs.slides[2]  # 削除後 index=2 が元の参考文献スライド

for shape in slide_ref.shapes:
    if shape.has_text_frame and '参考文献' in shape.text_frame.text:
        tf = shape.text_frame
        txBody = tf._txBody
        for p in txBody.findall(qn('a:p')):
            txBody.remove(p)
        new_refs = [
            ("〇参考文献", True, False),
            ("", False, False),
            ("OM 3-4 飛行実施計画 3-4-2 ⑦ Conflict Zone", False, False),
            ("OM S-3-11 Conflict Zoneおよびその周囲での運航について（Eff. 20260401）★NEW", True, True),
            ("OR OMR　２．Conflict Zone", False, False),
            ("運航関連業務管理規則　Conflict Zoneに係る飛行空域運用要領", False, False),
            ("Operation Update No.26003　Operations of Flights Considering the Situation in the Middle East　★NEW", True, True),
            ("Operation Update No.26005　Explanation of OM S-3-11 'Operations over and near Conflict Zones'　★NEW", True, True),
            ("IFALPA Flying into and Over Conflict Zones", False, False),
        ]
        for text, bold, is_new in new_refs:
            p_elem = etree.SubElement(txBody, qn('a:p'))
            r_elem = etree.SubElement(p_elem, qn('a:r'))
            rPr = etree.SubElement(r_elem, qn('a:rPr'),
                                   attrib={'lang':'ja-JP','altLang':'en-US'})
            if bold: rPr.set('b','1')
            if is_new:
                sf = etree.SubElement(rPr, qn('a:solidFill'))
                etree.SubElement(sf, qn('a:srgbClr'), attrib={'val':'CC0000'})
            t_elem = etree.SubElement(r_elem, qn('a:t'))
            t_elem.text = text
        break

# =========================================
# PHASE 7: スライド順序を整列
# 現在: [0=Title, 1=Category, 2=Ref, 3=RiskA, 4=ASR]
# 目標: [0=Title, 1=Category, 2=RiskA, 3=ASR, 4=Ref]
# =========================================
def move_slide(prs, old, new):
    sldIdLst = prs.slides._sldIdLst
    elem = sldIdLst[old]
    sldIdLst.remove(elem)
    sldIdLst.insert(new, elem)

move_slide(prs, 2, 4)  # Ref を末尾へ

# =========================================
# 保存
# =========================================
prs.save(DST)
print(f"\n保存完了: {DST}")
print(f"最終スライド数: {len(prs.slides)}枚\n")
for i, slide in enumerate(prs.slides):
    texts = [s.text_frame.paragraphs[0].text[:45].strip()
             for s in slide.shapes if s.has_text_frame
             if s.text_frame.paragraphs[0].text.strip()]
    print(f"  Slide {i+1}: {' | '.join(texts[:3])}")
