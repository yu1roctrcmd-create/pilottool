#!/usr/bin/env python3
"""NCA井戸端会議 Vol.2 Rev.1 PowerPoint generator"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ======== カラー定義 ========
C_BLUE       = RGBColor(0x00, 0x70, 0xC0)  # タイトル青
C_DARKBLUE   = RGBColor(0x2E, 0x75, 0xB6)  # テーブルヘッダー
C_LIGHTBLUE  = RGBColor(0xBD, 0xD7, 0xEE)  # テーブル項目列
C_PALE_BLUE  = RGBColor(0xEB, 0xF3, 0xFB)  # 背景等
C_YELLOW     = RGBColor(0xFF, 0xD7, 0x00)  # Pilot B吹き出し
C_GREEN      = RGBColor(0x70, 0xAD, 0x47)  # Pilot A吹き出し
C_ORANGE     = RGBColor(0xFF, 0xA0, 0x7A)  # Pilot C吹き出し
C_RED_HEAD   = RGBColor(0xC0, 0x00, 0x00)  # セクション見出し
C_HIGHLIGHT  = RGBColor(0xFF, 0xE6, 0x99)  # 比較表ハイライト
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK      = RGBColor(0x00, 0x00, 0x00)
C_GRAY       = RGBColor(0xF0, 0xF0, 0xF0)
C_DARKGRAY   = RGBColor(0x44, 0x44, 0x44)
C_BOX_BG     = RGBColor(0xF8, 0xF8, 0xF8)
C_BOX_BORDER = RGBColor(0x4A, 0x86, 0xC8)
C_EVEN_ROW   = RGBColor(0xDE, 0xEA, 0xF1)

# ======== プレゼンテーション設定 ========
prs = Presentation()
prs.slide_width  = Inches(13.33)   # ワイド 16:9
prs.slide_height = Inches(7.5)

SL_W = prs.slide_width
SL_H = prs.slide_height

blank_layout = prs.slide_layouts[6]  # 完全ブランク

# ======== ヘルパー関数 ========

def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, text, x, y, w, h,
                font_size=Pt(11), bold=False, color=C_BLACK,
                align=PP_ALIGN.LEFT, wrap=True,
                fill_color=None, line_color=None,
                v_anchor=None, line_width=Pt(0.75)):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    if v_anchor:
        tf.vertical_anchor = v_anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = 'Meiryo UI'
    # fill / border
    sp = txBox._element
    spPr = sp.find(qn('p:spPr'))
    if fill_color:
        solidFill = etree.SubElement(spPr, qn('a:solidFill'))
        srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgb.set('val', str(fill_color))
    if line_color:
        ln = etree.SubElement(spPr, qn('a:ln'))
        ln.set('w', str(int(line_width)))
        solidFill2 = etree.SubElement(ln, qn('a:solidFill'))
        srgb2 = etree.SubElement(solidFill2, qn('a:srgbClr'))
        srgb2.set('val', str(line_color))
    return txBox

def add_multiline_textbox(slide, lines, x, y, w, h,
                           font_size=Pt(10), bold=False, color=C_BLACK,
                           align=PP_ALIGN.LEFT,
                           fill_color=None, line_color=None,
                           title_line=None, line_width=Pt(1.5)):
    """複数行テキストボックス（規程ボックス用）"""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    sp = txBox._element
    spPr = sp.find(qn('p:spPr'))
    if fill_color:
        solidFill = etree.SubElement(spPr, qn('a:solidFill'))
        srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgb.set('val', str(fill_color))
    if line_color:
        ln_el = etree.SubElement(spPr, qn('a:ln'))
        ln_el.set('w', str(int(line_width)))
        solidFill2 = etree.SubElement(ln_el, qn('a:solidFill'))
        srgb2 = etree.SubElement(solidFill2, qn('a:srgbClr'))
        srgb2.set('val', str(line_color))
    first = True
    for i, line in enumerate(lines):
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = 'Meiryo UI'
        if title_line is not None and i == 0:
            run.font.size = Pt(font_size.pt + 0.5) if hasattr(font_size, 'pt') else font_size
            run.font.bold = True
            run.font.color.rgb = C_DARKBLUE
        else:
            run.font.size = font_size
            run.font.bold = bold
            run.font.color.rgb = color
    return txBox

def bubble(slide, speaker, text, x, y, w, h, bg_color=C_YELLOW, text_color=C_BLACK, font_size=Pt(10.5)):
    """吹き出し風テキストボックス（角丸矩形で近似）"""
    # 背景矩形
    shape = slide.shapes.add_shape(5, x, y, w, h)  # 5=rounded rectangle
    shape.adjustments[0] = 0.05
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
    shape.line.width = Pt(0.75)
    tf = shape.text_frame
    tf.word_wrap = True
    # スピーカーラベル
    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = speaker
    r0.font.size = Pt(8.5)
    r0.font.bold = True
    r0.font.color.rgb = C_DARKGRAY
    r0.font.name = 'Meiryo UI'
    # テキスト本文（改行で複数段落）
    for line in text.split('\n'):
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.size = font_size
        r.font.color.rgb = text_color
        r.font.name = 'Meiryo UI'
    return shape

def add_slide_header(slide, title, subtitle='', bg_color=C_PALE_BLUE):
    """スライド共通ヘッダー帯"""
    add_rect(slide, Inches(0), Inches(0), SL_W, Inches(0.85), fill_color=bg_color, line_color=None)
    # NCA小タイトル
    tb_sub = slide.shapes.add_textbox(Inches(0.2), Inches(0.0), Inches(6), Inches(0.35))
    tf = tb_sub.text_frame
    r = tf.paragraphs[0].add_run()
    r.text = 'SHIBAYAMA SMALL TOWN TALK  |  NCA 井戸端会議  Vol.2 Rev.1'
    r.font.size = Pt(8)
    r.font.color.rgb = C_BLUE
    r.font.bold = True
    r.font.name = 'Meiryo UI'
    # メインタイトル
    tb = slide.shapes.add_textbox(Inches(0.2), Inches(0.3), Inches(11), Inches(0.6))
    tf2 = tb.text_frame
    p = tf2.paragraphs[0]
    r2 = p.add_run()
    r2.text = title
    r2.font.size = Pt(18)
    r2.font.bold = True
    r2.font.color.rgb = C_RED_HEAD
    r2.font.name = 'Meiryo UI'
    # リファレンス（右端）
    tb_ref = slide.shapes.add_textbox(Inches(10.5), Inches(0.3), Inches(2.7), Inches(0.5))
    tf3 = tb_ref.text_frame
    tf3.paragraphs[0].alignment = PP_ALIGN.RIGHT
    r3 = tf3.paragraphs[0].add_run()
    r3.text = 'TYOOB 18-022-Rev1  |  Mar 2, 2026'
    r3.font.size = Pt(7.5)
    r3.font.color.rgb = C_DARKGRAY
    r3.font.name = 'Meiryo UI'

def add_slide_footer(slide, page_num, total=6):
    """スライド共通フッター"""
    add_rect(slide, Inches(0), Inches(7.15), SL_W, Inches(0.35), fill_color=C_DARKBLUE, line_color=None)
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(7.15), Inches(12), Inches(0.35))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = f'Reference#: TYOOB 18-022-Rev1   |   Issued by: TYOOBKZ   |   Prepared by: Hara  yosuke.hara@nca.aero   |   Page {page_num}/{total}'
    r.font.size = Pt(7.5)
    r.font.color.rgb = C_WHITE
    r.font.name = 'Meiryo UI'

def regulation_box(slide, title, lines, x, y, w, h, font_size=Pt(9.5)):
    add_multiline_textbox(slide, [title] + lines, x, y, w, h,
                          font_size=font_size,
                          fill_color=C_BOX_BG,
                          line_color=C_BOX_BORDER,
                          title_line=0,
                          line_width=Pt(1.5))

# ======================================================
#  SLIDE 1 ― タイトルスライド
# ======================================================
slide1 = prs.slides.add_slide(blank_layout)

# 背景グラデーション風
add_rect(slide1, 0, 0, SL_W, SL_H, fill_color=C_PALE_BLUE, line_color=None)
add_rect(slide1, 0, 0, SL_W, Inches(2.8), fill_color=C_DARKBLUE, line_color=None)

# メインタイトル
tb = slide1.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1.1))
tf = tb.text_frame
r = tf.paragraphs[0].add_run()
r.text = 'SHIBAYAMA SMALL TOWN TALK'
r.font.size = Pt(36)
r.font.bold = True
r.font.color.rgb = C_WHITE
r.font.name = 'Arial Black'

tb2 = slide1.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(0.7))
tf2 = tb2.text_frame
r2 = tf2.paragraphs[0].add_run()
r2.text = 'NCA 井戸端会議'
r2.font.size = Pt(28)
r2.font.bold = True
r2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xAA)
r2.font.name = 'Meiryo UI'

# Vol表示
tb3 = slide1.shapes.add_textbox(Inches(1), Inches(2.3), Inches(5), Inches(0.5))
tf3 = tb3.text_frame
r3 = tf3.paragraphs[0].add_run()
r3.text = 'Vol. 2  Rev.1　　　　　　　　　　　　　　　　　　March 2, 2026'
r3.font.size = Pt(13)
r3.font.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)
r3.font.name = 'Meiryo UI'

# サブタイトル（メインテーマ）
add_rect(slide1, Inches(0.8), Inches(3.2), Inches(11.7), Inches(1.1),
         fill_color=C_WHITE, line_color=C_BLUE)
tb4 = slide1.shapes.add_textbox(Inches(1.0), Inches(3.3), Inches(11.3), Inches(0.9))
tf4 = tb4.text_frame
tf4.paragraphs[0].alignment = PP_ALIGN.CENTER
r4 = tf4.paragraphs[0].add_run()
r4.text = 'アプローチ中のスピードコントロールについて（総集編）'
r4.font.size = Pt(22)
r4.font.bold = True
r4.font.color.rgb = C_RED_HEAD
r4.font.name = 'Meiryo UI'

# カバー内容目次
lines_idx = [
    '  ①  日本（AIP JAPAN ENR 1.8.6）',
    '  ②  香港（AIP HONG KONG / ICAO Doc 4444）',
    '  ③  FAA（FAA Order JO 7110.65）',
    '  ④  3地域まとめ比較',
]
tb5 = slide1.shapes.add_textbox(Inches(2.5), Inches(4.5), Inches(8.3), Inches(2.5))
tf5 = tb5.text_frame
tf5.word_wrap = True
for i, line in enumerate(lines_idx):
    if i == 0:
        p = tf5.paragraphs[0]
    else:
        p = tf5.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = line
    r.font.size = Pt(14)
    r.font.color.rgb = C_DARKBLUE
    r.font.name = 'Meiryo UI'

add_slide_footer(slide1, 1)

# ======================================================
#  SLIDE 2 ― シナリオ導入 ＋ 日本のルール
# ======================================================
slide2 = prs.slides.add_slide(blank_layout)
add_rect(slide2, 0, 0, SL_W, SL_H, fill_color=C_WHITE, line_color=None)
add_slide_header(slide2, '① 日本のルール（AIP JAPAN ENR 1.8.6）')
add_slide_footer(slide2, 2)

# シナリオ吹き出し（ATC交信）
regulation_box(slide2, '【シナリオ：Bさんが香港へアプローチ中】',
    ['NCA203  Reduce Speed 160kts',
     'NCA203  Turn Left Direct RIVER  Cleared ILS RWY 25R APP',
     'NCA203  Contact Tower 118.2'],
    Inches(0.3), Inches(1.0), Inches(6.2), Inches(1.3), font_size=Pt(10))

# Pilot B吹き出し（左）
bubble(slide2, '【Pilot B】',
    'あら？ Approach Clearanceが出たけど\n'
    '速度指示はCancelになったの？\n'
    '日本では自動で終了するはずだけど…',
    Inches(0.3), Inches(2.45), Inches(5.3), Inches(1.3),
    bg_color=C_YELLOW)

# Pilot A吹き出し（右）
bubble(slide2, '【Pilot A】',
    'ほな、日本の規程から\nおさらいしてみよかー！',
    Inches(7.7), Inches(2.45), Inches(5.3), Inches(1.1),
    bg_color=C_GREEN)

# AIP JAPAN 規程ボックス
regulation_box(slide2, 'AIP JAPAN  ENR 1.8.6',
    ['次に掲げる場合は、それ以前に指示されていた速度調整は 自動的に終了する。',
     'a)  待機が指示された場合',
     'b)  「Climb via SID」又は「Descend via STAR」が指示された場合',
     'c)  進入許可が発出され、管制官から速度調整に係る指示が新しく又は繰り返して行われなかった場合  ★',
     'd)  レーダー進入において接地点から5マイルの地点または最終降下開始点のうち、いずれか接地点から遠い方の地点を通過したのち',
     'e)  速度を維持すべき地点が明示されたのち当該地点を通過した場合'],
    Inches(0.3), Inches(3.75), Inches(8.3), Inches(2.85), font_size=Pt(9.5))

# まとめ吹き出し（右下）
bubble(slide2, '【Pilot B まとめ】',
    'c)より、Approach Clearance発出時点で\n速度調整は 自動的に終了！\n→ Pilot\'s Discretionで減速できる。\n（日本独自のルール）',
    Inches(8.8), Inches(4.0), Inches(4.3), Inches(2.3),
    bg_color=C_YELLOW)

# ======================================================
#  SLIDE 3 ― 香港（ICAO Doc 4444）
# ======================================================
slide3 = prs.slides.add_slide(blank_layout)
add_rect(slide3, 0, 0, SL_W, SL_H, fill_color=C_WHITE, line_color=None)
add_slide_header(slide3, '② 香港のルール（AIP HONG KONG / ICAO Doc 4444）')
add_slide_footer(slide3, 3)

# AIP HKG規程
regulation_box(slide3, 'AIP HONG KONG  ENR 1 GENERAL RULES AND PROCEDURES',
    ['1.1  The air traffic rules and procedures applicable to air traffic within the',
     '      Hong Kong FIR conform to: Annex 2 and Annex 11, the Air Navigation (Hong Kong) Order 1995;',
     '      ICAO Doc 4444 PANS/ATM, and the Regional Supplementary Procedures MID/ASIA Region,',
     '      except for the differences listed in GEN 1.7.'],
    Inches(0.3), Inches(1.0), Inches(12.7), Inches(1.5), font_size=Pt(9.5))

# Pilot B
bubble(slide3, '【Pilot B】',
    'ふむふむ、香港はICAO Doc 4444に従うのね。\nICAO Differencesを見たところ、速度調整に関する相違点はなさそう。\nDec 4444にはなんて書いてあるのかしら？',
    Inches(0.3), Inches(2.6), Inches(6.5), Inches(1.3),
    bg_color=C_YELLOW)

# Pilot A
bubble(slide3, '【Pilot A】',
    '勉強熱心やなー。\nほれ、これや。',
    Inches(7.5), Inches(2.6), Inches(5.5), Inches(1.1),
    bg_color=C_GREEN)

# ICAO Doc 4444
regulation_box(slide3, 'AIR TRAFFIC MANAGEMENT (DOC 4444)',
    ['4.6.1.2  Speed control instructions shall remain in effect unless explicitly',
     '              cancelled or amended by the controller.',
     '',
     '4.6.3.7  Speed control should not be applied to aircraft after passing a point',
     '              7 km (4 NM) from the threshold on final approach.',
     '',
     'NOTE:  The flight crew has a requirement to fly a stabilized approach typically by',
     '            5 km (3 NM) from the threshold (Doc 8168, PANS-OPS refers)'],
    Inches(0.3), Inches(3.85), Inches(7.8), Inches(2.7), font_size=Pt(9.5))

# まとめ吹き出し（右）
bubble(slide3, '【Pilot B まとめ】',
    '香港では Approach Clearanceが出ても\nSpeed ControlはCancelにならない！\n\n管制官が明示的にCancelするまで有効。\nThresholdから4NM以内への速度指示は不可。\n\n→ 日本との大きな違い！',
    Inches(8.3), Inches(3.85), Inches(4.8), Inches(2.7),
    bg_color=C_YELLOW)

# ======================================================
#  SLIDE 4 ― FAA（JO 7110.65）
# ======================================================
slide4 = prs.slides.add_slide(blank_layout)
add_rect(slide4, 0, 0, SL_W, SL_H, fill_color=C_WHITE, line_color=None)
add_slide_header(slide4, '③ FAA（米国）のルール（FAA Order JO 7110.65）')
add_slide_footer(slide4, 4)

# Pilot C登場
bubble(slide4, '【Pilot C】',
    'ひょっこりはん！ FAA編、俺に任せてや！',
    Inches(8.5), Inches(1.0), Inches(4.5), Inches(0.9),
    bg_color=C_ORANGE)

bubble(slide4, '【Pilot A】',
    '待ってたで、Pilot C！ FAAはどうなんや？',
    Inches(0.3), Inches(1.0), Inches(5.5), Inches(0.9),
    bg_color=C_GREEN)

# FAA規程ボックス
regulation_box(slide4, 'FAA Order JO 7110.65  Section 5-7-1  SPEED ADJUSTMENT',
    ['a.  Apply speed adjustment procedures to achieve or maintain required or desired spacing.',
     '',
     'b.  Do not assign speed adjustments to aircraft:',
     '      1.  At or inside the final approach fix (FAF) on final, or',
     '      2.  A point 5 miles from the runway,',
     '      whichever is closer to the runway.',
     '',
     'c.  Speed adjustment instructions shall remain in effect until explicitly cancelled or',
     '      modified by the controller.',
     '      ※ Approach Clearanceの発出のみでは速度指示はCancelされない。',
     '         管制官が明示的に「Cancel Speed Restrictions」等と指示した場合のみ終了する。'],
    Inches(0.3), Inches(2.1), Inches(7.8), Inches(3.7), font_size=Pt(9.5))

# まとめ吹き出し（右）
bubble(slide4, '【Pilot C まとめ】',
    'FAAもApproach Clearanceだけでは\nSpeed ControlはCancelにならない。\n管制官の明示的なCancelが必要！\n\n【ICAOとの違い】\nFAAは Runway から5NM以内\nまたはFAFの内側（近い方）が対象外。\n\nICAOはThresholdから4NM以内。\n→ FAFの位置によってはICAOより\n  早く制限対象外になることも。',
    Inches(8.3), Inches(2.1), Inches(4.8), Inches(4.0),
    bg_color=C_ORANGE)

# Pilot B
bubble(slide4, '【Pilot B】',
    'ILSのFAFはだいたい5〜7NM付近。\nFAAはFAFの存在がポイントになるのね！',
    Inches(0.3), Inches(5.9), Inches(6.5), Inches(1.0),
    bg_color=C_YELLOW)

# ======================================================
#  SLIDE 5 ― まとめ比較表
# ======================================================
slide5 = prs.slides.add_slide(blank_layout)
add_rect(slide5, 0, 0, SL_W, SL_H, fill_color=C_WHITE, line_color=None)
add_slide_header(slide5, '④ まとめ：3地域のスピードコントロールルール比較')
add_slide_footer(slide5, 5)

# Pilot B 導入
bubble(slide5, '【Pilot B】',
    '3地域を整理してみましょう！',
    Inches(0.3), Inches(0.95), Inches(5.0), Inches(0.75),
    bg_color=C_YELLOW)

# ======== 比較テーブル ========
rows = 5
cols = 4
left   = Inches(0.3)
top    = Inches(1.85)
width  = Inches(12.7)
height = Inches(3.5)

table = slide5.shapes.add_table(rows, cols, left, top, width, height).table

# 列幅
table.columns[0].width = Inches(3.5)
table.columns[1].width = Inches(3.0)
table.columns[2].width = Inches(3.1)
table.columns[3].width = Inches(3.1)

def set_cell(cell, text, bg_rgb, font_size=Pt(10), bold=False, font_color=C_BLACK, align=PP_ALIGN.CENTER):
    cell.text = ''
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = font_color
    run.font.name = 'Meiryo UI'
    # 背景色
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
    srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgb.set('val', str(bg_rgb))

# ヘッダー行
headers = [
    ('項目', C_DARKBLUE),
    ('日本  (AIP JAPAN)', C_DARKBLUE),
    ('香港  (ICAO Doc 4444)', C_DARKBLUE),
    ('FAA  (JO 7110.65)', C_DARKBLUE),
]
for col_idx, (text, bg) in enumerate(headers):
    set_cell(table.cell(0, col_idx), text, bg,
             font_size=Pt(11), bold=True, font_color=C_WHITE)

# データ行
data = [
    ['Approach Clearanceで\n速度指示が自動終了',
     '○\n（自動終了）',
     '✗\n（自動終了しない）',
     '✗\n（自動終了しない）'],
    ['速度指示の\n終了タイミング',
     'Approach Clearance\n発出時（c項）',
     '管制官が明示的に\nCancelした時',
     '管制官が明示的に\nCancelした時'],
    ['速度管理の\n対象外となる地点',
     '接地点から5NM\nまたは最終降下開始点\n（遠い方）',
     'Thresholdから\n4NM以内',
     'Runwayから5NM以内\nまたはFAFの内側\n（近い方）'],
    ['根拠規程',
     'AIP JAPAN\nENR 1.8.6',
     'ICAO Doc 4444\n4.6.1.2 / 4.6.3.7',
     'FAA Order JO 7110.65\nSection 5-7-1'],
]

for row_idx, row_data in enumerate(data):
    actual_row = row_idx + 1
    for col_idx, text in enumerate(row_data):
        bg = C_LIGHTBLUE if col_idx == 0 else (C_EVEN_ROW if row_idx % 2 == 0 else C_WHITE)
        # 日本の「自動終了○」を黄色ハイライト
        if col_idx == 1 and row_idx == 0:
            bg = C_HIGHLIGHT
        bold = (col_idx == 0) or (col_idx == 1 and row_idx == 0)
        font_color = C_RED_HEAD if (col_idx == 1 and row_idx == 0) else C_BLACK
        set_cell(table.cell(actual_row, col_idx), text, bg,
                 font_size=Pt(9.5), bold=bold, font_color=font_color)

# まとめコメント
bubble(slide5, '【Pilot A】',
    '日本は独自ルール！ Approach Clearanceで自動終了。\nHKGもFAAも明示的なCancelが必要。\n海外フライトでは必ずその国のAIPを確認せなあかんな！',
    Inches(0.3), Inches(5.5), Inches(6.3), Inches(1.5),
    bg_color=C_GREEN)

bubble(slide5, '【Pilot B】',
    'Approach Clearance後も速度指示は生きていると意識して、\nしっかり確認しながら飛ぶことが大切ね！',
    Inches(6.8), Inches(5.5), Inches(6.2), Inches(1.5),
    bg_color=C_YELLOW)

# ======================================================
#  SLIDE 6 ― エンディング
# ======================================================
slide6 = prs.slides.add_slide(blank_layout)
add_rect(slide6, 0, 0, SL_W, SL_H, fill_color=C_DARKBLUE, line_color=None)

# テキスト
tb = slide6.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.3), Inches(1.0))
tf = tb.text_frame
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
r = tf.paragraphs[0].add_run()
r.text = '好評につき、次回につづく！'
r.font.size = Pt(32)
r.font.bold = True
r.font.color.rgb = C_WHITE
r.font.name = 'Meiryo UI'

bubble(slide6, '【Pilot C】',
    '俺の出番、これだけ…\nひょっこりはん！',
    Inches(4.5), Inches(3.5), Inches(4.3), Inches(1.3),
    bg_color=C_ORANGE)

bubble(slide6, '【Pilot A】',
    'もうええわ、おおきに！\nでも他の国も調べなあかんな。',
    Inches(0.5), Inches(5.2), Inches(5.5), Inches(1.3),
    bg_color=C_GREEN)

bubble(slide6, '【Pilot B】',
    '次回もお楽しみに！',
    Inches(7.3), Inches(5.2), Inches(5.5), Inches(1.1),
    bg_color=C_YELLOW)

tb2 = slide6.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.5))
tf2 = tb2.text_frame
tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
r2 = tf2.paragraphs[0].add_run()
r2.text = 'Reference#: TYOOB 18-022-Rev1  |  Issued by: TYOOBKZ  |  Prepared by: Hara yosuke.hara@nca.aero  |  March 2, 2026'
r2.font.size = Pt(7.5)
r2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xFF)
r2.font.name = 'Meiryo UI'

# ======== SAVE ========
output_path = '/Users/yuichiromori/Desktop/NCAフォルダ/NCA井戸端会議_VOL.2_Rev1.pptx'
prs.save(output_path)
print(f'✓ PowerPoint saved: {output_path}')
