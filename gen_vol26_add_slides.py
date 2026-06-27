#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
VOL_26 Rev1草案に以下のスライドを追加:
  - リスク評価方法（Category 1/2/3の決定プロセス）
  - ASR報告義務（新規追加）
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

SRC = "/Users/yuichiromori/Desktop/Conflict Zones周辺の飛行 改定版（2026.04.01）JP_Rev1草案.pptx"
DST = "/Users/yuichiromori/Desktop/Conflict Zones周辺の飛行 改定版（2026.04.01）JP_Rev1草案.pptx"

prs = Presentation(SRC)
blank_layout = prs.slide_layouts[6]

# =========================================
# ヘルパー関数
# =========================================
def table_cell_fill(cell, hex_color):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for existing in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(existing)
    solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', hex_color)

def add_textbox(slide, text, left, top, width, height,
                font_size=10, bold=False, color=(0,0,0), align=PP_ALIGN.LEFT,
                bg_color=None, border_color=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    if bg_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = RGBColor(*bg_color)
    if border_color:
        txBox.line.color.rgb = RGBColor(*border_color)
        txBox.line.width = Pt(0.5)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return txBox

def add_multiline_textbox(slide, lines, left, top, width, height,
                          font_size=9, bg_color=None, border_color=None):
    """複数行テキストボックス。各行は(text, bold, color_rgb)のタプル"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    if bg_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = RGBColor(*bg_color)
    if border_color:
        txBox.line.color.rgb = RGBColor(*border_color)
        txBox.line.width = Pt(0.5)
    for i, (text, bold, color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(*color)
    return txBox

def add_header(slide, title_text):
    # 背景矩形（濃紺）
    bg = slide.shapes.add_shape(1, Inches(0.39), Inches(0.27), Inches(12.26), Inches(2.50))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(26, 54, 93)
    bg.line.fill.background()

    # タイトルテキスト群
    for left, top, w, h, txt, sz, bold in [
        (0.42, 0.36, 9.23, 0.75, "SHIBAYAMA SMALL TOWN TALK", 22, True),
        (3.5,  1.2,  5.50, 0.45, "NCA 井戸端会議", 14, True),
        (0.81, 2.0,  8.53, 0.35, "Vol. 26  Rev.1          Apr 01, 2026", 10, False),
    ]:
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
        p = tb.text_frame.paragraphs[0]
        if txt == "NCA 井戸端会議":
            p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.bold = bold
        r.font.color.rgb = RGBColor(255,255,255) if bold else RGBColor(200,220,255)

    # スライドタイトル（白背景の帯）
    title_bg = slide.shapes.add_shape(1, Inches(0.39), Inches(2.62), Inches(12.26), Inches(0.72))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    title_bg.line.fill.background()
    tb_t = slide.shapes.add_textbox(Inches(0.81), Inches(2.65), Inches(11.5), Inches(0.6))
    p_t = tb_t.text_frame.paragraphs[0]
    p_t.alignment = PP_ALIGN.CENTER
    r_t = p_t.add_run()
    r_t.text = title_text
    r_t.font.size = Pt(14)
    r_t.font.bold = True
    r_t.font.color.rgb = RGBColor(26, 54, 93)

def set_tbl_cell(tbl, row, col, text, fill_hex=None, bold=False, font_size=9,
                 font_color=(30,30,30), align=PP_ALIGN.LEFT):
    cell = tbl.cell(row, col)
    cell.text = text
    if fill_hex:
        table_cell_fill(cell, fill_hex)
    for para in cell.text_frame.paragraphs:
        para.alignment = align
        for run in para.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = RGBColor(*font_color)


# =========================================
# NEW Slide A: リスク評価方法
# =========================================
slide_a = prs.slides.add_slide(blank_layout)
add_header(slide_a, "カテゴリー評価方法 ─ リスクアセスメントの手順")

# 導入
add_textbox(slide_a,
    "カテゴリーは、ICAO Doc.10084 に基づき「T（脅威）× C（結果の重大性）× V（脆弱性）」の"
    "3要素をそれぞれ Rate 1〜3 で評価し、合計値でリスクレベル（High/Medium/Low）を決定。"
    "その結果をカテゴリー1〜3 に当てはめるよ。",
    Inches(0.8), Inches(3.55), Inches(11.5), Inches(0.55),
    font_size=9.5, color=(30,30,30))

# ---- 3要素ブロック ----
factors = [
    ("T　Threat（脅威）― 当該空域に存在する軍事的・技術的な危険性",
     [("Rate 3（高い）", True, (140,0,0)),
      ("　過去数年以内に実際に軍事的行動や攻撃が発生、またはその能力・意図・計画の明確な証拠がある。", False, (30,30,30)),
      ("Rate 2（中程度）", True, (140,80,0)),
      ("　比較的最近に、短期・中期の攻撃計画や数量的な意図を示す事例や証拠がある。", False, (30,30,30)),
      ("Rate 1（低い）", True, (0,100,0)),
      ("　最近の事例がなく、攻撃や攻撃計画等の兆候もない。", False, (30,30,30))],
     (255,215,215), (150,80,80)),
    ("C　Consequence（結果の重大性）― 機体および航空環境への影響の大きさ",
     [("Rate 3（重大）", True, (140,0,0)),
      ("　機体の損失、または管制・交通流への深刻な影響（交通流の麻痺等）が発生する。", False, (30,30,30)),
      ("Rate 2（中程度）", True, (140,80,0)),
      ("　機体の重篤な損傷、または管制・交通流への大きな影響（通過の大幅な遅延等）が発生する。", False, (30,30,30)),
      ("Rate 1（軽微）", True, (0,100,0)),
      ("　機体の部分的な損傷、または管制・交通流への部分的な影響が発生する。", False, (30,30,30))],
     (255,245,210), (160,130,0)),
    ("V　Vulnerability（脆弱性）― リスク緩和措置の実施状況",
     [("Rate 3（高い）", True, (140,0,0)),
      ("　リスクの緩和措置（リスク回避のための制限等）が実施されていない。", False, (30,30,30)),
      ("Rate 2（中程度）", True, (140,80,0)),
      ("　緩和措置は存在するが、現時点で効果的な対策の設定が困難である。", False, (30,30,30)),
      ("Rate 1（低い）", True, (0,100,0)),
      ("　効果的と認められるリスクの緩和措置が実施されている。", False, (30,30,30))],
     (225,245,220), (60,130,60)),
]

y = Inches(4.2)
for title, body_lines, bg, border in factors:
    # タイトルバー
    r_,g_,b_ = [max(0,c-60) for c in bg]
    title_bg = slide_a.shapes.add_shape(1, Inches(0.8), y, Inches(11.5), Inches(0.32))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = RGBColor(r_,g_,b_)
    title_bg.line.fill.background()
    tb_t = slide_a.shapes.add_textbox(Inches(0.85), y, Inches(11.4), Inches(0.32))
    p_t = tb_t.text_frame.paragraphs[0]
    r_t = p_t.add_run()
    r_t.text = title
    r_t.font.size = Pt(9.5)
    r_t.font.bold = True
    r_t.font.color.rgb = RGBColor(255,255,255)
    tb_t.text_frame.margin_top = Inches(0.02)
    y += Inches(0.32)

    # 本文
    bh = Inches(0.85)
    body_bg = slide_a.shapes.add_shape(1, Inches(0.8), y, Inches(11.5), bh)
    body_bg.fill.solid()
    body_bg.fill.fore_color.rgb = RGBColor(*bg)
    body_bg.line.color.rgb = RGBColor(*border)
    body_bg.line.width = Pt(0.5)
    add_multiline_textbox(slide_a, body_lines, Inches(0.85), y, Inches(11.3), bh,
                          font_size=8.5)
    y += bh + Inches(0.08)

# ---- リスクレベル → カテゴリー 決定表 ----
add_textbox(slide_a, "▼ T+C+V 合計値によるカテゴリー決定",
    Inches(0.8), y + Inches(0.05), Inches(11.5), Inches(0.35),
    font_size=10, bold=True, color=(26,54,93))
y += Inches(0.42)

tbl = slide_a.shapes.add_table(4, 5, Inches(0.8), y, Inches(11.5), Inches(1.5)).table
# カラム幅調整
tbl.columns[0].width = Inches(1.2)
tbl.columns[1].width = Inches(1.3)
tbl.columns[2].width = Inches(2.0)
tbl.columns[3].width = Inches(3.0)
tbl.columns[4].width = Inches(4.0)

for i, h in enumerate(["Risk Level", "T+C+V合計", "Category", "制限内容", "運航上の条件"]):
    set_tbl_cell(tbl, 0, i, h, "2E5FA3", bold=True, font_color=(255,255,255), align=PP_ALIGN.CENTER)

rows = [
    ("High",   "7〜9", "Category 1\n(Prohibited)",   "飛行禁止",    "リスク抑制のための対策を\n講じる必要がある",        "FFD7D7"),
    ("Medium", "4〜6", "Category 2\n(Restricted)",   "飛行制限",    "既にリスク抑制措置が含まれていることを確認\n（記録に含まれない要素も考慮）", "FFF2CC"),
    ("Low",    "1〜3", "Category 3\n(Information)",  "通常運航可",  "具体的な措置等は不要",                         "E2EFDA"),
]
for r_idx, (lvl, ttl, cat, limit, cond, color) in enumerate(rows):
    for c_idx, text in enumerate([lvl, ttl, cat, limit, cond]):
        set_tbl_cell(tbl, r_idx+1, c_idx, text, color,
                     align=PP_ALIGN.CENTER if c_idx < 3 else PP_ALIGN.LEFT)

add_textbox(slide_a,
    "※ カテゴリーの最終決定では T/C/V 合計に加え、社的状況・保険の適用・その他の緩和要素も勘案します。\n"
    "   評価結果は Risk Assessment Record に記録し、状況変化時は再評価を実施。（OR 運航関連業務管理規則 Appendix参照）",
    Inches(0.8), y + Inches(1.55), Inches(11.5), Inches(0.5),
    font_size=8, color=(100,100,100))


# =========================================
# NEW Slide B: 飛行後の報告（ASR）
# =========================================
slide_b = prs.slides.add_slide(blank_layout)
add_header(slide_b, "飛行後の報告 ── ASR 報告対象の新規追加")

add_textbox(slide_b,
    "OM S-3-11 の新設に伴い、Air Safety Report（ASR）の報告対象に Conflict Zone 関連事象が新たに追加されました。",
    Inches(0.8), Inches(3.55), Inches(11.5), Inches(0.45),
    font_size=10, bold=True, color=(180,0,0))

# ASR対象の表
tbl2 = slide_b.shapes.add_table(4, 3, Inches(0.8), Inches(4.1), Inches(11.5), Inches(1.7)).table
tbl2.columns[0].width = Inches(0.5)
tbl2.columns[1].width = Inches(8.0)
tbl2.columns[2].width = Inches(3.0)

for i, h in enumerate(["#", "ASR報告の対象となる事象", "関連規程"]):
    set_tbl_cell(tbl2, 0, i, h, "2E5FA3", bold=True, font_color=(255,255,255), align=PP_ALIGN.CENTER)

asr_items = [
    ("①", "悪天回避等によりConflict Zone（計画経路から50NM以内）へ入域した、または入域を検討した場合", "OM S-3-11 / OR 22章"),
    ("②", "GPS Jamming / Spoofing が発生した、または疑われた場合", "OM S-3-11 / OR 22章"),
    ("③", "Conflict Zone 飛行に係るその他の特異事象", "OM S-3-11"),
]
fills = ["F2F2F2", "FFFFFF", "F2F2F2"]
for r_idx, ((no, event, ref), fill) in enumerate(zip(asr_items, fills)):
    set_tbl_cell(tbl2, r_idx+1, 0, no,    fill, align=PP_ALIGN.CENTER)
    set_tbl_cell(tbl2, r_idx+1, 1, event, fill)
    set_tbl_cell(tbl2, r_idx+1, 2, ref,   fill)

# 不測の事態の対応
add_textbox(slide_b,
    "▼ 不測の事態で Conflict Zone への入域が避けられない場合の対処措置",
    Inches(0.8), Inches(5.95), Inches(11.5), Inches(0.38),
    font_size=10, bold=True, color=(26,54,93))

add_multiline_textbox(slide_b, [
    ("・ トランスポンダ 7700 のセット", False, (30,30,30)),
    ("・ 遭難通信の発信", False, (30,30,30)),
    ("・ 121.5 MHz での一方送信", False, (30,30,30)),
    ("・ 夜間のエクステリアライト点灯（通常 MEL 範囲内では機体識別への影響なし）", False, (30,30,30)),
    ("※ 緊急事態または安全上やむを得ない入域は直ちに危険をもたらすものではありません。安全最優先で対処すること。",
     True, (160,0,0)),
], Inches(0.8), Inches(6.38), Inches(11.5), Inches(1.4),
   font_size=9.5, bg_color=(255,245,220), border_color=(200,160,80))

add_textbox(slide_b,
    "（参照：OM S-3-11 3.5節「不測の事態への対応」/ OR 運航関連業務管理規則 22章）",
    Inches(0.8), Inches(7.85), Inches(11.5), Inches(0.3),
    font_size=8, color=(120,120,120))

# まとめの一言
add_multiline_textbox(slide_b, [
    ("飛行後の報告（ASR）は、会社としての安全情報の蓄積と将来的なリスク評価の精度向上に役立つよ。", False, (30,30,30)),
    ("Conflict Zone 関連の経験は積極的に報告してね！", True, (26,54,93)),
], Inches(0.8), Inches(8.3), Inches(11.5), Inches(0.6),
   font_size=9.5, bg_color=(230,240,255), border_color=(100,140,200))


# =========================================
# スライドを適切な位置に移動
# 現在: [0=Title, 1=Category, 2=Ref, 3=RiskA, 4=ASR]
# 目標: [0=Title, 1=Category, 2=RiskA, 3=ASR, 4=Ref]
# =========================================
def move_slide(prs, old_idx, new_idx):
    sldIdLst = prs.slides._sldIdLst
    slide_elem = sldIdLst[old_idx]
    sldIdLst.remove(slide_elem)
    sldIdLst.insert(new_idx, slide_elem)

# Ref（index=2）を末尾（index=4）へ
move_slide(prs, 2, 4)

# =========================================
# 保存
# =========================================
prs.save(DST)
print(f"保存完了: {DST}")
print(f"最終スライド数: {len(prs.slides)}枚")
for i, slide in enumerate(prs.slides):
    texts = [s.text_frame.paragraphs[0].text[:40].strip()
             for s in slide.shapes if s.has_text_frame
             if s.text_frame.paragraphs[0].text.strip()]
    print(f"  Slide {i+1}: {' | '.join(texts[:3])}")
